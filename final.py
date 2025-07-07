import streamlit as st
import validators
import os
import base64
from fpdf import FPDF
from langchain.prompts import PromptTemplate
from langchain.chains.summarize import load_summarize_chain
from langchain.document_loaders import PyPDFLoader, YoutubeLoader, WebBaseLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_groq import ChatGroq
from deep_translator import GoogleTranslator
import mysql.connector
import sqlite3
import re
import graphviz
import math
from gtts import gTTS  
from pptx import Presentation  
import tempfile 

# --- SQLite for Login/Register ---
conn_sqlite = sqlite3.connect("users.db", check_same_thread=False)
cursor_sqlite = conn_sqlite.cursor()

cursor_sqlite.execute("""
CREATE TABLE IF NOT EXISTS users (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    username TEXT UNIQUE NOT NULL,
    password TEXT NOT NULL
)""")
conn_sqlite.commit()

def add_user(username, password):
    try:
        cursor_sqlite.execute("INSERT INTO users (username, password) VALUES (?, ?)", (username, password))
        conn_sqlite.commit()
        return True
    except sqlite3.IntegrityError:
        return False

def user_exists(username):
    cursor_sqlite.execute("SELECT * FROM users WHERE username = ?", (username,))
    return cursor_sqlite.fetchone() is not None

def verify_user(username, password):
    cursor_sqlite.execute("SELECT * FROM users WHERE username = ? AND password = ?", (username, password))
    return cursor_sqlite.fetchone() is not None

# --- MySQL for summary history ---
def save_summary(user_key, source_type, source_value, summary):
    conn = mysql.connector.connect(host="localhost", user="root", password="1234", database="user_summary_history")
    cursor = conn.cursor()
    cursor.execute("INSERT INTO summaries (user_key, source_type, source_value, summary) VALUES (%s, %s, %s, %s)",
                   (user_key, source_type, source_value, summary))
    conn.commit()
    cursor.close()
    conn.close()

def get_summary_history(user_key):
    conn = mysql.connector.connect(host="localhost", user="root", password="1234", database="user_summary_history")
    cursor = conn.cursor(dictionary=True)
    cursor.execute("SELECT id, source_type, source_value, summary, timestamp FROM summaries WHERE user_key = %s ORDER BY timestamp DESC", (user_key,))
    rows = cursor.fetchall()
    cursor.close()
    conn.close()
    return rows

def delete_summary(summary_id):
    conn = mysql.connector.connect(host="localhost", user="root", password="1234", database="user_summary_history")
    cursor = conn.cursor()
    cursor.execute("DELETE FROM summaries WHERE id = %s", (summary_id,))
    conn.commit()
    cursor.close()
    conn.close()

def generate_presentation(summary_text, llm):
    """Convert summary to a colorful, attractive PowerPoint with 6-7 slides"""
    try:
        # Generate structured content from the summary
        prompt = f"""Transform this content into a professional PowerPoint presentation with 6-7 slides:
        {summary_text}
        
        Rules:
        1. Create 6-7 slides covering key aspects
        2. Don't copy exact text - rephrase for presentation
        3. Include:
           - Title slide with topic
           - Introduction/Overview
           - 3-4 key points (1 slide each)
           - Case study/example (if applicable)
           - Conclusion slide
        4. For each slide provide:
           - Slide title
           - 3-5 concise bullet points
           - Suggested visual theme (chart/diagram/image)
        
        Format exactly like this:
        Slide 1: [Title] | [Theme]
        - Content point 1
        - Content point 2
        
        Slide 2: [Title] | [Theme]
        - Content point 1
        - Content point 2"""
        
        response = llm.invoke(prompt)
        content = response.content if hasattr(response, 'content') else str(response)
        
        # Create presentation
        prs = Presentation()
        
        # Slide layouts
        title_slide_layout = prs.slide_layouts[0]
        section_header_layout = prs.slide_layouts[1]
        content_layout = prs.slide_layouts[5]  # Blank layout for custom designs
        
        # Color scheme
        colors = {
            'title_bg': (0, 32, 96),        # Dark blue
            'content_bg': (227, 239, 255),  # Light blue
            'accent1': (255, 102, 0),       # Orange
            'accent2': (0, 176, 80)         # Green
        }
        
        # Parse content and create slides
        current_slide = None
        for line in content.split('\n'):
            line = line.strip()
            if line.startswith('Slide'):
                # Add previous slide if exists
                if current_slide:
                    prs.slides.add_slide(current_slide)
                
                # Create new slide
                parts = line.split('|')
                title = parts[0].split(':')[1].strip()
                theme = parts[1].strip() if len(parts) > 1 else "Default"
                
                if "Title" in line:
                    current_slide = prs.slides.add_slide(title_slide_layout)
                    background = current_slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = colors['title_bg']
                    
                    title_shape = current_slide.shapes.title
                    title_shape.text = title
                    title_shape.text_frame.paragraphs[0].font.color.rgb = (255, 255, 255)  # White text
                    
                    subtitle = current_slide.placeholders[1]
                    subtitle.text = "AI-Generated Presentation"
                else:
                    current_slide = prs.slides.add_slide(content_layout)
                    
                    # Add title
                    title_box = current_slide.shapes.add_textbox(left=0.5*914400, top=0.2*914400, 
                                                              width=6*914400, height=914400)
                    tf = title_box.text_frame
                    p = tf.add_paragraph()
                    p.text = title
                    p.font.size = 32
                    p.font.color.rgb = colors['title_bg']
                    p.font.bold = True
                    
                    # Add theme indicator
                    theme_box = current_slide.shapes.add_textbox(left=7*914400, top=7*914400,
                                                               width=2*914400, height=0.5*914400)
                    tf = theme_box.text_frame
                    p = tf.add_paragraph()
                    p.text = f"Theme: {theme}"
                    p.font.size = 12
                    p.font.color.rgb = (100, 100, 100)
                    
            elif line.startswith('-'):
                if current_slide and not isinstance(current_slide, int):
                    # Add content to current slide
                    left = 1*914400
                    top = 1.5*914400 + (len(current_slide.shapes.placeholders)*0.5*914400)
                    
                    content_box = current_slide.shapes.add_textbox(left=left, top=top,
                                                                 width=7*914400, height=0.8*914400)
                    tf = content_box.text_frame
                    p = tf.add_paragraph()
                    p.text = "• " + line[1:].strip()
                    p.font.size = 24
                    p.font.color.rgb = (0, 0, 0)
                    p.space_after = 0
        
        # Add final slide if exists
        if current_slide and not isinstance(current_slide, int):
            prs.slides.add_slide(current_slide)
        
        # Add conclusion slide
        conclusion = prs.slides.add_slide(section_header_layout)
        conclusion.shapes.title.text = "Key Takeaways"
        conclusion.placeholders[1].text = "1. First key point\n2. Second key point\n3. Third key point"
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(temp_file.name)
        return temp_file.name
        
    except Exception as e:
        st.error(f"PPT Generation Error: {str(e)}")
        return None
def generate_quiz(summary_text, llm):
    """Generate and administer an interactive quiz from text with enhanced feedback"""
    try:
        prompt = f"""
        Create exactly 3 multiple choice quiz questions based on this text:
        {summary_text}
        
        Format each question exactly like this:
        Q) Question text here?
        A) Option 1
        B) Option 2
        C) Option 3
        D) Option 4
        Answer: X) Correct option letter (A, B, C, or D)
        Explanation: Brief explanation of why this is the correct answer
        
        Rules:
        1. Each question should test key information from the text
        2. Options should be plausible but only one correct answer
        3. Questions should cover different aspects of the content
        4. Include clear explanations for correct answers
        """
        response = llm.invoke(prompt)
        quiz_content = response.content if hasattr(response, 'content') else str(response)
        
        # Parse the quiz questions
        questions = []
        current_question = {}
        
        for line in quiz_content.split('\n'):
            line = line.strip()
            if line.startswith('Q)'):
                if current_question:  # Save previous question if exists
                    questions.append(current_question)
                current_question = {
                    'question': line[2:].strip(),
                    'options': [],
                    'answer': None,
                    'explanation': None
                }
            elif line.startswith(('A)', 'B)', 'C)', 'D)')):
                current_question['options'].append(line[2:].strip())
            elif line.startswith('Answer:'):
                current_question['answer'] = line.split(')')[0][-1].upper()
            elif line.startswith('Explanation:'):
                current_question['explanation'] = line[12:].strip()
        
        if current_question:  # Add the last question
            questions.append(current_question)
        
        # Limit to exactly 3 questions
        questions = questions[:3]
        
        # Initialize session state
        if 'quiz_data' not in st.session_state:
            st.session_state.quiz_data = {
                'questions': questions,
                'answers': {},
                'submitted': False,
                'score': 0
            }
        
        st.subheader("📝 Quiz")
        st.write("Answer all 3 questions based on the summary, then click Submit:")
        
        # Display questions
        for i, q in enumerate(st.session_state.quiz_data['questions']):
            # Get current answer or None if not answered yet
            current_answer = st.session_state.quiz_data['answers'].get(i)
            
            # Display question
            st.markdown(f"**Q{i+1}) {q['question']}**")
            
            # Create radio buttons for options
            user_answer = st.radio(
                "Select your answer:",
                options=q['options'],
                key=f"quiz_q_{i}",
                index=q['options'].index(current_answer) if current_answer else None
            )
            
            # Store answer in session state
            if user_answer and (i not in st.session_state.quiz_data['answers'] or 
                              st.session_state.quiz_data['answers'][i] != user_answer):
                st.session_state.quiz_data['answers'][i] = user_answer
        
        # Submit button
        if st.button("Submit Answers") and not st.session_state.quiz_data['submitted']:
            if len(st.session_state.quiz_data['answers']) < 3:
                st.warning("Please answer all questions before submitting!")
            else:
                # Calculate score
                score = 0
                feedback = []
                for i, q in enumerate(st.session_state.quiz_data['questions']):
                    correct_answer = q['options'][ord(q['answer']) - ord('A')]
                    is_correct = st.session_state.quiz_data['answers'][i] == correct_answer
                    if is_correct:
                        score += 1
                    
                    feedback.append({
                        'question': q['question'],
                        'user_answer': st.session_state.quiz_data['answers'][i],
                        'correct_answer': correct_answer,
                        'is_correct': is_correct,
                        'explanation': q['explanation']
                    })
                
                st.session_state.quiz_data['submitted'] = True
                st.session_state.quiz_data['score'] = score
                st.session_state.quiz_data['feedback'] = feedback
                st.rerun()
        
        # Show results if submitted
        if st.session_state.quiz_data.get('submitted', False):
            st.markdown("---")
            st.subheader(f"🎯 Quiz Results: {st.session_state.quiz_data['score']}/3")
            
            # Performance summary
            performance = ""
            score = st.session_state.quiz_data['score']
            if score == 3:
                performance = "🎉 Excellent! You got all questions right!"
            elif score == 2:
                performance = "👍 Good job! You got most questions right."
            elif score == 1:
                performance = "📚 Not bad! Review the material and try again."
            else:
                performance = "🔄 Let's review the content and try again."
            
            st.markdown(f"### {performance}")
            
            # Detailed feedback for all questions
            st.markdown("### 📝 Review Your Answers:")
            for i, feedback in enumerate(st.session_state.quiz_data['feedback']):
                st.markdown(f"**Q{i+1}) {feedback['question']}**")
                st.markdown(f"Your answer: {feedback['user_answer']}")
                
                if not feedback['is_correct']:
                    st.markdown(f"**Correct answer:** {feedback['correct_answer']}")
                
                st.markdown(f"**Explanation:** {feedback['explanation']}")
                st.markdown("---")
            
            if st.button("🔄 Retake Quiz"):
                st.session_state.quiz_data = {
                    'questions': questions,
                    'answers': {},
                    'submitted': False,
                    'score': 0
                }
                st.rerun()
        
        return None
        
    except Exception as e:
        st.error(f"Quiz generation failed: {str(e)}")
        return None

def generate_audiobook(summary_text, llm):
    """Generate a fresh 2-minute audio gist of the topic"""
    try:
        # Generate a fresh, engaging audio script
        prompt = f"""Create a concise 2-minute (150-200 word) audio script about this topic:
        {summary_text}
        
        Rules:
        1. DON'T read the summary verbatim
        2. Create fresh, engaging content that covers:
           - What the topic is about
           - Why it matters
           - Key insights
           - Practical applications
        3. Use conversational language
        4. Keep sentences short
        5. Maximum 200 words
        
        Example format:
        "Today we'll explore... 
        The key thing to understand is... What makes this important is...
        You can apply this by..." """
        
        response = llm.invoke(prompt)
        script = response.content if hasattr(response, 'content') else str(response)
        
        # Further limit to ~200 words (approximately 1.5-2 minutes)
        words = script.split()[:200]
        final_script = ' '.join(words)
        
        # Convert to speech with natural pacing
        tts = gTTS(
            text=final_script,
            lang='en',
            slow=False,
            lang_check=False
        )
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3")
        tts.save(temp_file.name)
        return temp_file.name
        
    except Exception as e:
        st.error(f"Audio generation error: {str(e)}")
        return None
# --- Text Cleaning Utilities ---
def clean_text(text):
    """Clean text of problematic unicode characters"""
    replacements = {
        '\u2013': '-',  # en dash
        '\u2014': '-',  # em dash
        '\u2018': "'",  # left single quote
        '\u2019': "'",  # right single quote
        '\u201c': '"',  # left double quote
        '\u201d': '"',  # right double quote
    }
    for k, v in replacements.items():
        text = text.replace(k, v)
    return text

# --- Download Utilities ---
def generate_pdf(summary_text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    # Use built-in Helvetica font which supports basic characters
    pdf.set_font("Helvetica", size=12)
    
    # Clean text before adding to PDF
    clean_text = summary_text.encode('latin-1', 'replace').decode('latin-1')
    
    for line in clean_text.split('\n'):
        pdf.multi_cell(0, 10, line)
    pdf_path = "summary.pdf"
    pdf.output(pdf_path)
    return pdf_path

def get_download_link(file_path, label, file_type):
    with open(file_path, "rb") as f:
        bytes_data = f.read()
    b64 = base64.b64encode(bytes_data).decode()
    href = f'<a href="data:application/{file_type};base64,{b64}" download="{os.path.basename(file_path)}">{label}</a>'
    return href

# --- Translation ---
def translate_text(text, target_lang):
    if target_lang == "en":
        return text
    try:
        return GoogleTranslator(source='auto', target=target_lang).translate(text)
    except:
        return text

# --- Password Validation ---
def is_valid_password(password):
    if len(password) < 6:
        return False
    if not re.search(r'[a-z]', password):
        return False
    if not re.search(r'[A-Z]', password):
        return False
    if not re.search(r'\d', password):
        return False
    return True

# --- Mindmap Generation ---
def generate_mindmap(summary_text, llm):
    """
    Generate a visual mindmap diagram with radial layout similar to the example image
    Returns graphviz Digraph object and text structure
    """
    try:
        # Enhanced prompt for better visual structure
        prompt = f"""Convert this content into a radial mindmap structure:
        {summary_text}
        
        Return in this exact format:
        # Central Topic
        - Branch 1
        - Branch 2
        * Sub-branch 2.1
        * Sub-branch 2.2
        - Branch 3
        
        Rules:
        1. Use # for central topic (single word if possible)
        2. Use - for main branches (1-2 words)
        3. Use * for sub-branches (1-3 words)
        4. Create 4-6 main branches
        5. Keep sub-branches concise"""
        
        response = llm.invoke(prompt)
        structure = response.content if hasattr(response, 'content') else str(response)
        
        # Create radial graph
        graph = graphviz.Digraph(
            graph_attr={
                'layout': 'neato',
                'overlap': 'false',
                'splines': 'true',
                'bgcolor': 'transparent',
                'center': 'true'
            }
        )
        
        # Color palette
        colors = {
            'central': '#FF6B6B',  # Red
            'branches': ['#4ECDC4', '#45B7D1', '#A5D8D1', '#7FB3D5'],  # Blues/Teals
            'sub': ['#FFE66D', '#FFD166', '#F7C1BB']  # Yellows/Pinks
        }
        
        # Parse structure and build graph
        central_node = None
        nodes = {}
        
        for line in structure.split('\n'):
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('#'):
                # Central node
                central_node = line[1:].strip().split()[0]  # Take first word only
                graph.node(central_node, 
                         shape='circle',
                         style='filled',
                         fillcolor=colors['central'],
                         fontcolor='white',
                         fontsize='24',
                         width='1',
                         height='1',
                         fixedsize='true')
                nodes[central_node] = central_node
                
            elif line.startswith('-'):
                # Main branch
                if not central_node:
                    continue
                branch = ' '.join(line[1:].strip().split()[:2])  # Max 2 words
                color = colors['branches'][len(nodes) % len(colors['branches'])]
                graph.node(branch,
                         shape='ellipse',
                         style='filled',
                         fillcolor=color,
                         fontcolor='white',
                         fontsize='18')
                graph.edge(central_node, branch,
                          penwidth='2',
                          color='#555555')
                nodes[branch] = branch
                
            elif line.startswith('*'):
                # Sub-branch
                if not nodes:
                    continue
                sub = ' '.join(line[1:].strip().split()[:3])  # Max 3 words
                last_branch = list(nodes.keys())[-1]
                color = colors['sub'][len(nodes) % len(colors['sub'])]
                graph.node(sub,
                         shape='box',
                         style='rounded,filled',
                         fillcolor=color,
                         fontsize='14',
                         margin='0.1,0.05')
                graph.edge(last_branch, sub,
                          penwidth='1.5',
                          color='#777777',
                          arrowsize='0.7')
                nodes[sub] = sub
        
        # Position nodes radially
        angle = 360 / max(1, len([n for n in nodes if n != central_node]))
        for i, node in enumerate(nodes):
            if node == central_node:
                graph.node(node, pos='0,0!')
            else:
                radius = 2 if node in nodes and node.lower() != central_node.lower() else 1
                x = radius * math.cos(math.radians(angle * i))
                y = radius * math.sin(math.radians(angle * i))
                graph.node(node, pos=f'{x},{y}!')
        
        return graph, structure
        
    except Exception as e:
        st.error(f"Mindmap generation error: {str(e)}")
        return None, None
# --- Chatbot for Summary ---
def summary_chatbot(summary_text, summary_id, llm):
    # Initialize chat history if not exists
    if f"chat_history_{summary_id}" not in st.session_state:
        st.session_state[f"chat_history_{summary_id}"] = []
    
    # Display chat interface outside the expander
    st.markdown("---")
    st.subheader("💬 Chat with this Summary")
    
    # Display chat messages
    for message in st.session_state[f"chat_history_{summary_id}"]:
        with st.chat_message(message["role"]):
            st.write(message["content"])
    
    # Handle user input
    if prompt := st.chat_input("Ask about this summary..."):
        # Add user message to chat history
        st.session_state[f"chat_history_{summary_id}"].append({"role": "user", "content": prompt})
        
        # Display user message immediately
        with st.chat_message("user"):
            st.write(prompt)
        
        try:
            # Generate assistant response
            context = f"""
            Summary Context:
            {summary_text}
            
            User Question: {prompt}
            
            Answer the question based only on the summary above.
            If you can't answer, say 'This information is not in the summary'.
            """
            
            response = llm.invoke(context)
            answer = response.content if hasattr(response, 'content') else str(response)
            
            # Add assistant response to chat history
            st.session_state[f"chat_history_{summary_id}"].append({"role": "assistant", "content": answer})
            
            # Display assistant response
            with st.chat_message("assistant"):
                st.write(answer)
                
        except Exception as e:
            st.error(f"Error generating response: {str(e)}")

# --- Streamlit Page Config ---
st.set_page_config(page_title="Smart Summarizer App", layout="centered")

# --- Session state init ---
for key in ["logged_in", "username", "page", "target_lang", "active_chat", "role"]:
    if key not in st.session_state:
        st.session_state[key] = "" if key == "username" else False if key == "logged_in" else "role_select" if key == "page" else "en" if key == "target_lang" else None
# --- Language Map & Selector ---
language_map = {"English": "en", "Hindi": "hi", "Telugu": "te", "Tamil": "ta", "Kannada": "kn", "French": "fr", "Spanish": "es", "German": "de"}
selected_language = st.sidebar.selectbox("🌐 Select Language", list(language_map.keys()), index=0)
st.session_state.target_lang = language_map[selected_language]
target_lang = st.session_state.target_lang

# --- Login Page ---
def login_page():
    st.title("🔐 " + translate_text("Login", target_lang))
    username = st.text_input(translate_text("Username", target_lang))
    password = st.text_input(translate_text("Password", target_lang), type="password")
    if st.button(translate_text("Login", target_lang)):
        if not user_exists(username):
            st.error(translate_text("Account does not exist. Please create an account.", target_lang))
        elif verify_user(username, password):
            st.success(translate_text("Login successful!", target_lang))
            st.session_state.logged_in = True
            st.session_state.username = username
            st.session_state.page = "main"
        else:
            st.error(translate_text("Invalid username or password.", target_lang))

    st.write("---")
    if st.button(translate_text("Go to Register", target_lang)):
        st.session_state.page = "register"
    if st.button(translate_text("Forgot Password?", target_lang)):
        st.session_state.page = "forgot_password"

# --- Forgot Password Page ---
def forgot_password_page():
    st.title("🔑 " + translate_text("Reset Password", target_lang))
    username = st.text_input(translate_text("Enter your username", target_lang))
    new_password = st.text_input(translate_text("Enter new password", target_lang), type="password")
    confirm_password = st.text_input(translate_text("Confirm new password", target_lang), type="password")

    if st.button(translate_text("Reset Password", target_lang)):
        if not user_exists(username):
            st.error(translate_text("Username not found.", target_lang))
        elif new_password != confirm_password:
            st.error(translate_text("Passwords do not match.", target_lang))
        elif not is_valid_password(new_password):
            st.error(translate_text("Password must be at least 6 characters long, with at least one uppercase, lowercase, and digit.", target_lang))
        else:
            cursor_sqlite.execute("UPDATE users SET password = ? WHERE username = ?", (new_password, username))
            conn_sqlite.commit()
            st.success(translate_text("Password reset successful! Please login.", target_lang))
            st.session_state.page = "login"
            st.rerun()

# --- Register Page ---
def register_page():
    st.title("📝 " + translate_text("Register", target_lang))
    username = st.text_input(translate_text("Choose Username", target_lang))
    password = st.text_input(translate_text("Choose Password", target_lang), type="password")
    password2 = st.text_input(translate_text("Confirm Password", target_lang), type="password")

    if st.button(translate_text("Register", target_lang)):
        if password != password2:
            st.error(translate_text("Passwords do not match.", target_lang))
        elif not is_valid_password(password):
            st.error(translate_text("Password must be at least 6 characters long, with at least one uppercase letter, one lowercase letter, and one digit.", target_lang))
        elif add_user(username, password):
            st.success(translate_text("Registration successful! Please login.", target_lang))
            st.session_state.page = "login"
        else:
            st.error(translate_text("Username already exists.", target_lang))

    st.write("---")
    if st.button(translate_text("Go to Login", target_lang)):
        st.session_state.page = "login"
def role_selection_page():
    st.title("👋 Welcome to Smart Summarizer")
    st.write("Please select your role to get customized summaries:")
    
    cols = st.columns(3)
    with cols[0]:
        if st.button("👩‍🏫 Teacher"):
            st.session_state.role = "teacher"
            st.session_state.page = "login"
            st.rerun()
        st.write("Formal tone with key insights")
    with cols[1]:
        if st.button("🧑‍🎓 Student"):
            st.session_state.role = "student"
            st.session_state.page = "login"
            st.rerun()
        st.write("Detailed explanations with examples")
    with cols[2]:
        if st.button("👦 Kid (Under 12)"):
            st.session_state.role = "kid"
            st.session_state.page = "login"
            st.rerun()
        st.write("Simple language with fun examples")
def role_selection_page():
    st.title("👋 Welcome to Smart Summarizer")
    st.write("Please select your role to get customized summaries:")
    
    cols = st.columns(3)
    with cols[0]:
        if st.button("👩‍🏫 Teacher"):
            st.session_state.role = "teacher"
            st.session_state.page = "login"
            st.rerun()
        st.write("Formal tone with key insights")
    with cols[1]:
        if st.button("🧑‍🎓 Student"):
            st.session_state.role = "student"
            st.session_state.page = "login"
            st.rerun()
        st.write("Detailed explanations with examples")
    with cols[2]:
        if st.button("👦 Kid (Under 12)"):
            st.session_state.role = "kid"
            st.session_state.page = "login"
            st.rerun()
        st.write("Simple language with fun examples")
# Replace the existing prompt_template with:
def get_role_prompt(role):
    role_instructions = {
        "teacher": "Focus on key insights, main concepts, and pedagogical value. Use formal academic tone.",
        "student": "Provide detailed explanations with examples. Include study tips and practical applications.",
        "kid": "Use simple language (5th grade level). Include fun examples and analogies. Keep sentences short."
    }
    return role_instructions.get(role, "")

prompt_template = PromptTemplate.from_template(f"""
You are a helpful assistant summarizing content for a {st.session_state.role}. 
{get_role_prompt(st.session_state.role)}

Content: {{text}}

Return the summary in markdown format with appropriate headings.
""")
# --- Main Page ---
def main_page():
    role_display = {
        "teacher": "👩‍🏫 Teacher Mode",
        "student": "🧑‍🎓 Student Mode", 
        "kid": "👦 Kid Mode"
    }
    st.sidebar.markdown(f"**Current Mode:** {role_display.get(st.session_state.role, '')}")
    
    # Add role switching option
    if st.sidebar.button("Switch Role"):
        st.session_state.page = "role_select"
        st.rerun()
    st.title("📘 " + translate_text("Smart Summarizer App", target_lang))
    st.sidebar.title(translate_text("Settings", target_lang))
    groq_api_key = st.sidebar.text_input(translate_text("Groq API Key", target_lang), type="password")

    option = st.selectbox(translate_text("Select content type:", target_lang), [
        translate_text("Select", target_lang),
        translate_text("Summarize Website/YouTube", target_lang),
        translate_text("Summarize PDF Textbook", target_lang)
    ])

    if groq_api_key:
        try:
            llm = ChatGroq(model="llama3-8b-8192", groq_api_key=groq_api_key, temperature=0.3)
        except Exception as e:
            st.error(translate_text(f"Error initializing Groq model: {e}", target_lang))
            return

        prompt_template = PromptTemplate.from_template("""
You are a helpful assistant summarizing the following content. Use headings (##) and subheadings (###) to organize the summary. Keep it concise and informative.

Content: {text}

Return the summary in markdown format.
""")
        chain = load_summarize_chain(llm, chain_type="stuff", prompt=prompt_template)

        if option == translate_text("Summarize Website/YouTube", target_lang):
            st.header("🌐 " + translate_text("Summarize from Website or YouTube", target_lang))
            url = st.text_input(translate_text("Enter a website or YouTube URL", target_lang))
            if validators.url(url) and st.button(translate_text("Generate Summary", target_lang)):
                try:
                    loader = YoutubeLoader.from_youtube_url(url) if "youtube" in url.lower() else WebBaseLoader(url)
                    documents = loader.load()
                    if not documents:
                        st.warning(translate_text("No content found at the URL.", target_lang))
                        return
                    docs = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100).split_documents(documents)[:3]
                    summary = "\n\n".join(chain.run([doc]) for doc in docs)
                    summary = clean_text(summary)  # Clean text before translation
                    if target_lang != "en":
                        summary = GoogleTranslator(source='auto', target=target_lang).translate(summary)
                    st.success(translate_text("Summary generated!", target_lang))
                    st.markdown(summary)
                    save_summary(st.session_state.username, "url", url, summary)
                    with open("summary.txt", "w", encoding="utf-8") as f:
                        f.write(summary)
                    pdf_path = generate_pdf(summary)
                    st.markdown(get_download_link("summary.txt", "📄 " + translate_text("Download TXT", target_lang), "octet-stream"), unsafe_allow_html=True)
                    st.markdown(get_download_link(pdf_path, "📝 " + translate_text("Download PDF", target_lang), "pdf"), unsafe_allow_html=True)
                except Exception as e:
                    st.error(translate_text(f"Error during URL summarization: {e}", target_lang))

        elif option == translate_text("Summarize PDF Textbook", target_lang):
            st.header("📄 " + translate_text("Summarize PDF Textbook", target_lang))
            uploaded_pdf = st.file_uploader(translate_text("Upload your PDF file", target_lang), type=["pdf"])
            if uploaded_pdf and st.button(translate_text("Generate Summary", target_lang)):
                try:
                    with open("temp_uploaded_file.pdf", "wb") as f:
                        f.write(uploaded_pdf.read())
                    loader = PyPDFLoader("temp_uploaded_file.pdf")
                    documents = loader.load()
                    docs = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100).split_documents(documents)[:3]
                    summary = "\n\n".join(chain.run([doc]) for doc in docs)
                    summary = clean_text(summary)  # Clean text before translation
                    if target_lang != "en":
                        summary = GoogleTranslator(source='auto', target=target_lang).translate(summary)
                    st.success(translate_text("Summary generated!", target_lang))
                    st.markdown(summary)
                    save_summary(st.session_state.username, "pdf", uploaded_pdf.name, summary)
                    with open("summary.txt", "w", encoding="utf-8") as f:
                        f.write(summary)
                    pdf_path = generate_pdf(summary)
                    st.markdown(get_download_link("summary.txt", "📄 " + translate_text("Download TXT", target_lang), "octet-stream"), unsafe_allow_html=True)
                    st.markdown(get_download_link(pdf_path, "📝 " + translate_text("Download PDF", target_lang), "pdf"), unsafe_allow_html=True)
                    os.remove("temp_uploaded_file.pdf")
                except Exception as e:
                    st.error(translate_text(f"Error during PDF summarization: {e}", target_lang))
    else:
        st.warning(translate_text("Please enter your Groq API key.", target_lang))

    st.write("---")
    st.subheader(translate_text("Your Summary History", target_lang))
    history = get_summary_history(st.session_state.username)
    if history:
        for item in history:
            with st.expander(f"[{item['source_type'].upper()}] {item['source_value']} - {item['timestamp']}"):
                st.markdown(item['summary'])
                
                # Action buttons for each summary
                col1, col2, col3, col4 = st.columns([1,1,1,1])
                
                with col1:
                    if st.button(f"🧠 Mindmap", key=f"mindmap_{item['id']}"):
                        st.session_state[f"show_mindmap_{item['id']}"] = True
                        st.session_state[f"show_chat_{item['id']}"] = False
                        st.session_state[f"show_quiz_{item['id']}"] = False
                
                with col2:
                    if st.button(f"💬 Chat", key=f"chat_{item['id']}"):
                        st.session_state[f"show_chat_{item['id']}"] = True
                        st.session_state[f"show_mindmap_{item['id']}"] = False
                        st.session_state[f"show_quiz_{item['id']}"] = False
                
                with col3:
                    if st.button(f"❓ Quiz", key=f"quiz_{item['id']}"):
                        st.session_state[f"show_quiz_{item['id']}"] = True
                        st.session_state[f"show_mindmap_{item['id']}"] = False
                        st.session_state[f"show_chat_{item['id']}"] = False
                
                with col4:
                    if st.button(translate_text("🗑️ Delete", target_lang), key=f"delete_{item['id']}"):
                        delete_summary(item['id'])
                        st.rerun()
                
                # Show mindmap if requested
                if st.session_state.get(f"show_mindmap_{item['id']}"):
                    st.subheader("🧠 Mindmap")
                    graph, structure = generate_mindmap(item['summary'], llm)
                    if graph:
                        st.graphviz_chart(graph, use_container_width=True)
                
                # Show chatbot if requested
                if st.session_state.get(f"show_chat_{item['id']}"):
                    summary_chatbot(item['summary'], item['id'], llm)
                
                # Show quiz if requested
                if st.session_state.get(f"show_quiz_{item['id']}"):
                    st.subheader("📝 Quiz")
                    if st.button("Generate New Quiz", key=f"new_quiz_{item['id']}"):
                        st.session_state.quiz_answers = {}
                        st.session_state.quiz_score = 0
                        st.rerun()
                    generate_quiz(item['summary'], llm)
                
                # Audio option - always visible when expanded
                st.markdown("---")
                st.subheader("🔊 Audio Summary")
                if st.button("▶️ Listen to Summary (2 min)", key=f"audio_{item['id']}"):
                    audio_file = generate_audiobook(item['summary'],llm)
                    if audio_file:
                        st.audio(audio_file)
                        # Clean up the temporary file after playback
                        try:
                            os.unlink(audio_file)
                        except:
                            pass
    else:
        st.info(translate_text("No summary history found.", target_lang))
    if st.sidebar.button(translate_text("Logout", target_lang)):
        st.session_state.logged_in = False
        st.session_state.username = ""
        st.session_state.page = "login"
        st.rerun()


# --- Routing ---
if st.session_state.page == "role_select":
    role_selection_page()
elif st.session_state.page == "login":
    login_page()
elif st.session_state.page == "register":
    register_page()
elif st.session_state.page == "forgot_password":
    forgot_password_page()
elif st.session_state.page == "main" and st.session_state.logged_in:
    main_page()
else:
    st.session_state.page = "role_select"
    role_selection_page()